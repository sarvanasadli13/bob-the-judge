"""Builds the IBM-styled Bob the Judge pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# IBM Carbon colours
IBM_BLUE    = RGBColor(0x0f, 0x62, 0xfe)
IBM_DARK    = RGBColor(0x16, 0x16, 0x16)
IBM_GREY90  = RGBColor(0x26, 0x26, 0x26)
IBM_GREY60  = RGBColor(0x52, 0x52, 0x52)
IBM_GREY30  = RGBColor(0x8d, 0x8d, 0x8d)
IBM_WHITE   = RGBColor(0xf4, 0xf4, 0xf4)
GREEN       = RGBColor(0x42, 0xbe, 0x65)
RED         = RGBColor(0xfa, 0x4d, 0x56)
AMBER       = RGBColor(0xf1, 0xc2, 0x1b)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def box(slide, l, t, w, h, colour, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h, size, colour, bold=False, align=PP_ALIGN.LEFT, italic=False, font="IBM Plex Sans"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = colour
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return txb


def ibm_bar(slide, subtitle=""):
    box(slide, 0, 0, 13.33, 0.65, IBM_BLUE)
    txt(slide, "IBM", 0.3, 0.1, 1.2, 0.5, 22, IBM_WHITE, bold=True, font="IBM Plex Sans")
    if subtitle:
        txt(slide, subtitle, 1.6, 0.15, 10, 0.4, 13, RGBColor(0xc6,0xc6,0xc6), font="IBM Plex Sans")


def slide_divider(slide, colour=IBM_BLUE):
    box(slide, 0.5, 1.05, 12.33, 0.04, colour)


# ── SLIDE 1 — Title ────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1, IBM_DARK)
ibm_bar(s1, "Bob the Judge  ·  IBM Bob Hackathon 2026")
box(s1, 0, 0.65, 0.08, 6.85, IBM_BLUE)

txt(s1, "Bob the Judge", 0.5, 1.2, 10, 1.4, 64, IBM_WHITE, bold=True)
txt(s1, "Migration Cutover Decision Advisor", 0.5, 2.8, 10, 0.7, 28, IBM_GREY30)
txt(s1, "Powered by IBM Bob + IBM watsonx.ai Granite", 0.5, 3.6, 12, 0.5, 16, IBM_BLUE, bold=True)

box(s1, 0.5, 4.6, 12, 0.04, IBM_GREY60)
txt(s1, "Live demo:", 0.5, 4.75, 2.5, 0.4, 13, IBM_GREY30)
txt(s1, "bob-the-judge-production.up.railway.app", 2.0, 4.75, 11, 0.4, 13, GREEN, bold=True, font="IBM Plex Mono")

box(s1, 0.5, 5.3, 3.5, 0.06, IBM_BLUE)
txt(s1, "Bob writes the code.  Bob ships it.", 0.5, 5.5, 12, 0.5, 18, IBM_GREY30, italic=True)
txt(s1, "Bob the Judge tells Bob when.", 0.5, 5.95, 12, 0.5, 22, IBM_WHITE, bold=True)

txt(s1, "CyberFalcon Team  ·  IBM Bob Hackathon 2026", 0.5, 6.85, 12, 0.4, 12, IBM_GREY60, font="IBM Plex Mono")


# ── SLIDE 2 — Problem ─────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2, IBM_DARK)
ibm_bar(s2, "The Problem")
box(s2, 0, 0.65, 0.08, 6.85, RED)

txt(s2, "The Question Nobody Can Answer", 0.5, 0.9, 12, 0.8, 36, IBM_WHITE, bold=True)
slide_divider(s2, RED)

txt(s2, "When is it safe to cut over?", 0.5, 1.4, 12, 0.8, 28, RED, bold=True, italic=True)

stats = [
    ("70%", "of enterprise migrations\nstall in dual-run"),
    ("$200K", "per month running\ntwo systems in parallel"),
    ("0",     "tools that answer\nthe cutover question"),
]
for i, (num, label) in enumerate(stats):
    x = 0.5 + i * 4.2
    box(s2, x, 2.3, 3.7, 2.2, IBM_GREY90)
    box(s2, x, 2.3, 3.7, 0.06, RED)
    txt(s2, num,   x+0.2, 2.5,  3.3, 1.0, 52, RED,      bold=True)
    txt(s2, label, x+0.2, 3.5,  3.3, 0.9, 14, IBM_GREY30)

txt(s2, "CTOs make the cutover decision on gut feel. There is no objective, auditable answer.",
    0.5, 5.0, 12, 0.6, 16, IBM_GREY30, italic=True)


# ── SLIDE 3 — Solution ────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3, IBM_DARK)
ibm_bar(s3, "The Solution")
box(s3, 0, 0.65, 0.08, 6.85, IBM_BLUE)

txt(s3, "Bob the Judge", 0.5, 0.9, 12, 0.8, 40, IBM_WHITE, bold=True)
slide_divider(s3, IBM_BLUE)

points = [
    ("Function-level verdict",    "Per-function readiness score, not a system-level guess"),
    ("Real-time parity analysis", "Routes live traffic through both systems simultaneously"),
    ("Bob-native intelligence",   "Plan, Code, Ask, Orchestrator — Bob operates in all 4 modes"),
    ("Regulator-grade audit PDF", "One click exports a compliance sign-off document"),
]
for i, (title, desc) in enumerate(points):
    y = 1.5 + i * 1.2
    box(s3, 0.5, y, 0.06, 0.9, IBM_BLUE)
    txt(s3, title, 0.75, y,      11, 0.45, 18, IBM_WHITE, bold=True)
    txt(s3, desc,  0.75, y+0.42, 11, 0.45, 13, IBM_GREY30)

box(s3, 6.8, 1.5, 6.0, 5.2, IBM_GREY90)
box(s3, 6.8, 1.5, 6.0, 0.06, IBM_BLUE)
txt(s3, "VERDICT EXAMPLE", 6.9, 1.65, 5.8, 0.4, 10, IBM_GREY30, bold=True)

verdicts = [
    ("Domestic Wire",       "100", "SAFE TO CUT",  GREEN),
    ("Scheduled Payment",   "100", "SAFE TO CUT",  GREEN),
    ("International Wire",  "0",   "DO NOT CUT",   RED),
]
for i, (fn, score, verdict, colour) in enumerate(verdicts):
    y = 2.1 + i * 1.5
    box(s3, 6.9, y, 5.7, 1.3, IBM_DARK)
    box(s3, 6.9, y, 5.7, 0.05, colour)
    txt(s3, fn,     6.95, y+0.1,  3.5, 0.35, 11, IBM_GREY30)
    txt(s3, score,  6.95, y+0.4,  1.5, 0.7,  40, colour, bold=True)
    txt(s3, verdict,8.6,  y+0.55, 3.8, 0.4,  12, colour, bold=True)


# ── SLIDE 4 — How It Works ───────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4, IBM_DARK)
ibm_bar(s4, "How It Works")
box(s4, 0, 0.65, 0.08, 6.85, IBM_BLUE)

txt(s4, "Architecture", 0.5, 0.9, 12, 0.6, 36, IBM_WHITE, bold=True)
slide_divider(s4)

steps = [
    ("01", "Traffic Generator",  "Sends identical payment\ntransactions to both systems"),
    ("02", "Parity Engine",      "Compares fee, amount,\nstatus — flags divergences"),
    ("03", "Scoring Engine",     "Readiness score 0-100\nper payment function"),
    ("04", "IBM Bob",            "Plan + Code + Ask +\nOrchestrator modes"),
    ("05", "Audit PDF",          "Regulator-grade report\nwith sign-off page"),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.3 + i * 2.6
    box(s4, x, 1.6, 2.3, 3.8, IBM_GREY90)
    box(s4, x, 1.6, 2.3, 0.06, IBM_BLUE)
    txt(s4, num,   x+0.15, 1.75, 2.0, 0.6, 32, IBM_BLUE, bold=True, font="IBM Plex Mono")
    txt(s4, title, x+0.15, 2.45, 2.0, 0.5, 15, IBM_WHITE, bold=True)
    txt(s4, desc,  x+0.15, 3.0,  2.0, 1.2, 12, IBM_GREY30)
    if i < 4:
        txt(s4, ">", x+2.1, 2.9, 0.4, 0.5, 20, IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)

txt(s4, "Domain: Banking payment processing  ·  Legacy=COBOL-style, Modern=FastAPI  ·  Divergence by design for demo clarity",
    0.3, 5.8, 12.7, 0.4, 11, IBM_GREY30, italic=True)


# ── SLIDE 5 — Bob Integration ─────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5, IBM_DARK)
ibm_bar(s5, "IBM Bob Integration")
box(s5, 0, 0.65, 0.08, 6.85, IBM_BLUE)

txt(s5, "Bob Operates in All 4 Modes", 0.5, 0.9, 12, 0.6, 36, IBM_WHITE, bold=True)
slide_divider(s5)

modes = [
    ("Plan",         "#0f62fe", "Generates phased cutover\nstrategy with risk register\nand timeline"),
    ("Code",         "#0f62fe", "Produces ready-to-deploy\npatch for the highest-\nseverity divergence"),
    ("Ask",          "#0f62fe", "Explains each verdict in\nplain English with\nregulatory context"),
    ("Orchestrator", "#0f62fe", "Runs 4-agent pipeline:\nAnalyser > Risk-Scout >\nFix-Gen > Reporter"),
]
for i, (mode, colour, desc) in enumerate(modes):
    x = 0.4 + i * 3.2
    box(s5, x, 1.6, 2.9, 4.2, IBM_GREY90)
    box(s5, x, 1.6, 2.9, 0.5, IBM_BLUE)
    txt(s5, mode, x+0.15, 1.65, 2.6, 0.4, 16, IBM_WHITE, bold=True)
    txt(s5, desc, x+0.15, 2.3,  2.6, 2.0, 13, IBM_GREY30)

box(s5, 0.5, 6.0, 12.3, 0.7, IBM_GREY90)
box(s5, 0.5, 6.0, 0.06, 0.7, GREEN)
txt(s5, "● LIVE", 0.7, 6.07, 1.2, 0.4, 12, GREEN, bold=True, font="IBM Plex Mono")
txt(s5, "Powered by IBM Granite (ibm/granite-4-h-small) on watsonx.ai  ·  WATSONX-xxxx session IDs in every audit PDF",
    2.0, 6.1, 11, 0.4, 12, IBM_WHITE, italic=True)


# ── SLIDE 6 — Beyond Bob: watsonx.ai + Phase 1-4 by Bob ──────────────────────
s5b = prs.slides.add_slide(BLANK)
bg(s5b, IBM_DARK)
ibm_bar(s5b, "Beyond Bob  ·  watsonx.ai + Bob-Written Code")
box(s5b, 0, 0.65, 0.08, 6.85, IBM_BLUE)

txt(s5b, "The Full IBM Stack", 0.5, 0.9, 12, 0.6, 36, IBM_WHITE, bold=True)
slide_divider(s5b)

# Left column: watsonx.ai Granite
box(s5b, 0.5, 1.5, 6.0, 5.4, IBM_GREY90)
box(s5b, 0.5, 1.5, 6.0, 0.06, IBM_BLUE)
txt(s5b, "watsonx.ai", 0.7, 1.65, 5.5, 0.4, 11, IBM_BLUE, bold=True, font="IBM Plex Mono")
txt(s5b, "IBM Granite — Live", 0.7, 2.05, 5.5, 0.6, 22, IBM_WHITE, bold=True)
txt(s5b, "ibm/granite-4-h-small", 0.7, 2.7, 5.5, 0.4, 13, GREEN, font="IBM Plex Mono")
items_left = [
    "Real REST calls to us-south.ml.cloud.ibm.com",
    "Plan / Code / Ask / Orchestrator all routed through Granite",
    "WATSONX-xxxx session IDs in audit PDF chain-of-custody",
    "Graceful fallback to intelligent mock if API unavailable",
]
for i, it in enumerate(items_left):
    txt(s5b, "•", 0.7, 3.3 + i*0.55, 0.3, 0.4, 13, IBM_BLUE, bold=True)
    txt(s5b, it, 1.0, 3.3 + i*0.55, 5.2, 0.5, 12, IBM_GREY30)

# Right column: Bob-written enhancements
box(s5b, 6.8, 1.5, 6.0, 5.4, IBM_GREY90)
box(s5b, 6.8, 1.5, 6.0, 0.06, GREEN)
txt(s5b, "BUILT BY BOB", 7.0, 1.65, 5.5, 0.4, 11, GREEN, bold=True, font="IBM Plex Mono")
txt(s5b, "Phase 1–4 Modules", 7.0, 2.05, 5.5, 0.6, 22, IBM_WHITE, bold=True)
txt(s5b, "5 IBM Bob task session reports committed", 7.0, 2.7, 5.5, 0.4, 13, IBM_BLUE, italic=True)
items_right = [
    "Tenant calibration — 4 bank-tier profiles (FFIEC, Basel III)",
    "Confidence intervals — Wilson + Bootstrap + Bayesian",
    "Drift detection — CUSUM + EWMA on SQLite history",
    "Regulatory compliance — FFIEC, Basel III, PSD2, SOX 404",
    "48 pytest tests, 100% passing",
]
for i, it in enumerate(items_right):
    txt(s5b, "•", 7.0, 3.3 + i*0.5, 0.3, 0.4, 13, GREEN, bold=True)
    txt(s5b, it, 7.3, 3.3 + i*0.5, 5.2, 0.5, 12, IBM_GREY30)


# ── SLIDE 7 — Business Value ──────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6, IBM_DARK)
ibm_bar(s6, "Business Value")
box(s6, 0, 0.65, 0.08, 6.85, GREEN)

txt(s6, "Market & Impact", 0.5, 0.9, 12, 0.6, 36, IBM_WHITE, bold=True)
slide_divider(s6, GREEN)

metrics = [
    ("800B", "lines of COBOL\nstill in production"),
    ("$80B", "spent annually on\nmainframe modernization"),
    ("70%",  "migrations stall\nin dual-run phase"),
    ("$0",   "tools solving\nthe cutover decision"),
]
for i, (num, label) in enumerate(metrics):
    x = 0.4 + i * 3.2
    box(s6, x, 1.5, 2.9, 2.0, IBM_GREY90)
    box(s6, x, 1.5, 2.9, 0.06, GREEN)
    txt(s6, num,   x+0.2, 1.65, 2.5, 0.9, 40, GREEN, bold=True, font="IBM Plex Mono")
    txt(s6, label, x+0.2, 2.5,  2.5, 0.8, 13, IBM_GREY30)

box(s6, 0.5, 3.8, 12.3, 2.1, IBM_GREY90)
box(s6, 0.5, 3.8, 12.3, 0.06, GREEN)
advantages = [
    "Full IBM stack — Bob IDE built it, watsonx.ai Granite runs it in production",
    "Live deployment — judges can click the URL today, not just clone the repo",
    "Function-level granularity — competitors operate at workflow level",
    "Regulator-grade — FFIEC, Basel III, PSD2, SWIFT gpi citations are baked into the audit PDF",
]
txt(s6, "Why We Win", 0.7, 3.9, 5, 0.4, 14, GREEN, bold=True)
for i, adv in enumerate(advantages):
    txt(s6, f"  {adv}", 0.7, 4.35 + i*0.37, 12, 0.35, 12, IBM_GREY30)


# ── SLIDE 7 — Customer Voice (illustrative) ────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7, IBM_DARK)
ibm_bar(s7, "Customer Voice")
box(s7, 0, 0.65, 0.08, 6.85, IBM_BLUE)

txt(s7, "What CTOs Are Asking For", 0.5, 0.9, 12, 0.6, 36, IBM_WHITE, bold=True)
slide_divider(s7)

# Big quote block
box(s7, 0.8, 1.85, 11.7, 3.6, IBM_GREY90)
box(s7, 0.8, 1.85, 0.06, 3.6, IBM_BLUE)

txt(s7, "“", 1.0, 1.85, 1.5, 1.5, 80, IBM_BLUE, bold=True, italic=True, font="IBM Plex Serif")

quote_line_1 = "We've been stuck in dual-run for 14 months."
quote_line_2 = "Three consultants couldn't tell us when to cut over."
quote_line_3 = "Bob the Judge gave us per-function verdicts in 60 seconds —"
quote_line_4 = "and the audit PDF is exactly what our regulator asked for."

txt(s7, quote_line_1, 1.6, 2.05, 10.5, 0.5, 19, IBM_WHITE, italic=True)
txt(s7, quote_line_2, 1.6, 2.55, 10.5, 0.5, 19, IBM_WHITE, italic=True)
txt(s7, quote_line_3, 1.6, 3.15, 10.5, 0.5, 19, IBM_BLUE, italic=True, bold=True)
txt(s7, quote_line_4, 1.6, 3.65, 10.5, 0.5, 19, IBM_BLUE, italic=True, bold=True)

box(s7, 1.6, 4.45, 4.0, 0.04, IBM_GREY60)
txt(s7, "Sarah Chen, CTO", 1.6, 4.55, 6, 0.4, 14, IBM_WHITE, bold=True)
txt(s7, "Meridian Financial  ·  Mid-market US bank, $42B AUM", 1.6, 4.85, 8, 0.4, 12, IBM_GREY30)

txt(s7, "(Illustrative customer scenario based on observed pain points across COBOL banking modernization programs, 2025–2026.)",
    0.8, 5.7, 12, 0.5, 11, IBM_GREY60, italic=True)

# What this means
box(s7, 0.8, 6.3, 11.7, 0.7, IBM_GREY90)
box(s7, 0.8, 6.3, 0.06, 0.7, GREEN)
txt(s7, "The buyer is the CTO running a stalled migration. The price is one stalled month of dual-run infrastructure.",
    1.0, 6.45, 12, 0.4, 13, IBM_GREY30, italic=True)


# ── SLIDE 8 — Close ────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
bg(s8, IBM_DARK)
ibm_bar(s8, "Bob the Judge  ·  IBM Bob Hackathon 2026")
box(s8, 0, 0.65, 13.33, 6.85, IBM_DARK)
box(s8, 0, 0.65, 0.12, 6.85, IBM_BLUE)

txt(s8, "Bob writes the code.  Bob ships it.", 0.6, 2.0, 12, 0.9, 32, IBM_GREY30, italic=True)
txt(s8, "Bob the Judge tells Bob", 0.6, 2.85, 12, 0.9, 44, IBM_WHITE, bold=True)
txt(s8, "when to flip the switch.", 0.6, 3.65, 12, 0.9, 44, IBM_BLUE, bold=True)

box(s8, 0.6, 4.8, 4.0, 0.06, IBM_BLUE)
txt(s8, "The autonomous decision layer Bob was missing", 0.6, 5.0, 10, 0.5, 14, IBM_GREY30)
txt(s8, "Built by CyberFalcon Team  ·  IBM Bob Hackathon 2026", 0.6, 5.35, 12, 0.4, 13, IBM_BLUE, bold=True, font="IBM Plex Mono")

box(s8, 0.6, 5.9, 12.1, 0.7, IBM_GREY90)
box(s8, 0.6, 5.9, 0.06, 0.7, GREEN)
txt(s8, "● Try it live:", 0.8, 6.0, 2.5, 0.4, 13, GREEN, bold=True)
txt(s8, "bob-the-judge-production.up.railway.app", 2.5, 6.0, 11, 0.4, 14, IBM_WHITE, bold=True, font="IBM Plex Mono")
txt(s8, "github.com/sarvanasadli13/bob-the-judge  ·  5 IBM Bob session reports in /bob_sessions",
    0.8, 6.32, 12, 0.3, 11, IBM_GREY30)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "Bob_the_Judge_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
